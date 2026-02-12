#!/usr/bin/env python3
"""
PowerPoint Generator for Phishing Detection Project
Creates a 30-slide presentation with all sections
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

def create_presentation():
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Define colors
    PRIMARY_COLOR = RGBColor(0, 51, 102)  # Navy blue
    ACCENT_COLOR = RGBColor(0, 153, 204)  # Light blue
    TEXT_COLOR = RGBColor(32, 32, 32)     # Dark gray
    WHITE = RGBColor(255, 255, 255)
    
    def add_title_slide(prs, title, subtitle):
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Background shape
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = PRIMARY_COLOR
        shape.line.fill.background()
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        
        # Subtitle
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))
        tf = subtitle_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(28)
        p.font.color.rgb = ACCENT_COLOR
        p.alignment = PP_ALIGN.CENTER
        
        return slide
    
    def add_content_slide(prs, title, bullet_points, slide_num=None):
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Title bar
        title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
        title_bar.fill.solid()
        title_bar.fill.fore_color.rgb = PRIMARY_COLOR
        title_bar.line.fill.background()
        
        # Title text
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = WHITE
        
        # Slide number
        if slide_num:
            num_box = slide.shapes.add_textbox(Inches(12), Inches(0.4), Inches(1), Inches(0.5))
            tf = num_box.text_frame
            p = tf.paragraphs[0]
            p.text = str(slide_num)
            p.font.size = Pt(14)
            p.font.color.rgb = ACCENT_COLOR
        
        # Content
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.333), Inches(5.5))
        tf = content_box.text_frame
        tf.word_wrap = True
        
        for i, point in enumerate(bullet_points):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = "• " + point
            p.font.size = Pt(22)
            p.font.color.rgb = TEXT_COLOR
            p.space_after = Pt(12)
        
        return slide
    
    def add_two_column_slide(prs, title, left_title, left_points, right_title, right_points, slide_num=None):
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Title bar
        title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
        title_bar.fill.solid()
        title_bar.fill.fore_color.rgb = PRIMARY_COLOR
        title_bar.line.fill.background()
        
        # Title text
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = WHITE
        
        # Slide number
        if slide_num:
            num_box = slide.shapes.add_textbox(Inches(12), Inches(0.4), Inches(1), Inches(0.5))
            tf = num_box.text_frame
            p = tf.paragraphs[0]
            p.text = str(slide_num)
            p.font.size = Pt(14)
            p.font.color.rgb = ACCENT_COLOR
        
        # Left column title
        left_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(6), Inches(0.5))
        tf = left_title_box.text_frame
        p = tf.paragraphs[0]
        p.text = left_title
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = PRIMARY_COLOR
        
        # Left content
        left_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.1), Inches(6), Inches(5))
        tf = left_box.text_frame
        tf.word_wrap = True
        
        for i, point in enumerate(left_points):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = "• " + point
            p.font.size = Pt(18)
            p.font.color.rgb = TEXT_COLOR
            p.space_after = Pt(8)
        
        # Right column title
        right_title_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.5), Inches(6), Inches(0.5))
        tf = right_title_box.text_frame
        p = tf.paragraphs[0]
        p.text = right_title
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = PRIMARY_COLOR
        
        # Right content
        right_box = slide.shapes.add_textbox(Inches(6.8), Inches(2.1), Inches(6), Inches(5))
        tf = right_box.text_frame
        tf.word_wrap = True
        
        for i, point in enumerate(right_points):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = "• " + point
            p.font.size = Pt(18)
            p.font.color.rgb = TEXT_COLOR
            p.space_after = Pt(8)
        
        return slide
    
    def add_table_slide(prs, title, data, slide_num=None):
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Title bar
        title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
        title_bar.fill.solid()
        title_bar.fill.fore_color.rgb = PRIMARY_COLOR
        title_bar.line.fill.background()
        
        # Title text
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = WHITE
        
        # Slide number
        if slide_num:
            num_box = slide.shapes.add_textbox(Inches(12), Inches(0.4), Inches(1), Inches(0.5))
            tf = num_box.text_frame
            p = tf.paragraphs[0]
            p.text = str(slide_num)
            p.font.size = Pt(14)
            p.font.color.rgb = ACCENT_COLOR
        
        # Table
        rows = len(data)
        cols = len(data[0]) if rows > 0 else 0
        
        if rows > 0 and cols > 0:
            left = Inches(0.5)
            top = Inches(1.6)
            width = prs.slide_width - Inches(1)
            height = Inches(0.5) * rows
            
            table = slide.shapes.add_table(rows, cols, left, top, width, height).table
            
            # Set column widths (convert to integer value)
            col_width = int((prs.slide_width - Inches(1)) / cols)
            for i in range(cols):
                table.columns[i].width = col_width
            
            # Fill data
            for i, row_data in enumerate(data):
                for j, cell_text in enumerate(row_data):
                    cell = table.cell(i, j)
                    cell.text = cell_text
                    paragraph = cell.text_frame.paragraphs[0]
                    paragraph.font.size = Pt(14)
                    paragraph.font.color.rgb = TEXT_COLOR
                    
                    if i == 0:  # Header row
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = ACCENT_COLOR
                        paragraph.font.bold = True
                        paragraph.font.color.rgb = WHITE
        
        return slide
    
    # =========================================================================
    # SLIDE 1: TITLE SLIDE
    # =========================================================================
    add_title_slide(
        prs,
        "PHISHING GUARD V2.0",
        "AI-Powered Multimodal Phishing Detection System\nA Production-Grade IEEE-Level Final Year Project"
    )
    
    # =========================================================================
    # SLIDE 2: ABSTRACT - Overview
    # =========================================================================
    add_content_slide(
        prs,
        "ABSTRACT - Overview",
        [
            "Phishing Guard v2.0 is an advanced AI-based phishing detection system",
            "Innovative application of Multimodal Large Language Models (MLLMs)",
            "Transforms phishing detection into sophisticated text classification",
            "Analyzes multiple data modalities: URL, HTML, screenshots, and DOM features",
            "Achieves exceptional 99.8% F1 score",
            "Detects 4 distinct categories: LEGITIMATE, PHISHING, AI_GENERATED_PHISHING, PHISHING_KIT"
        ],
        2
    )
    
    # =========================================================================
    # SLIDE 3: ABSTRACT - Technical Approach
    # =========================================================================
    add_two_column_slide(
        prs,
        "ABSTRACT - Technical Approach",
        "Key Innovations",
        [
            "93 engineered features (+365% improvement)",
            "4-tier detection pipeline for resource efficiency",
            "Real-time protection via browser extension and REST API",
            "Enterprise-grade security with JWT authentication",
            "First production system for AI-generated phishing detection",
            "Qwen2.5-3B-Instruct with 4-bit quantization"
        ],
        "Impact & Applications",
        [
            "Personal cybersecurity protection",
            "Email gateway integration",
            "Enterprise API security",
            "Research and educational purposes",
            "50+ global and Indian brand protection",
            "Addresses emerging AI threat landscape"
        ],
        3
    )
    
    # =========================================================================
    # SLIDE 4: INTRODUCTION - The Digital Threat Landscape
    # =========================================================================
    add_content_slide(
        prs,
        "INTRODUCTION - The Digital Threat Landscape",
        [
            "5.3 billion internet users worldwide (2024)",
            "$8 trillion estimated global cost of cybercrime",
            "3.4 billion phishing emails sent daily",
            "$4.5 billion lost to phishing attacks in 2024 alone",
            "Phishing exploits human psychology, not just technical vulnerabilities",
            "Requires intelligent, adaptive detection systems"
        ],
        4
    )
    
    # =========================================================================
    # SLIDE 5: INTRODUCTION - Evolution of Phishing
    # =========================================================================
    add_two_column_slide(
        prs,
        "INTRODUCTION - Evolution of Phishing Attacks",
        "Generation 1: Traditional Phishing",
        [
            "Manual creation of deceptive websites",
            "Generic templates easily detectable",
            "Limited scalability",
            "Basic heuristics can identify"
        ],
        "Generation 2: Phishing Kits",
        [
            "Automated toolkits (Gophish, HiddenEye)",
            "Professional-grade deployment",
            "Scalable attack infrastructure",
            "Still detectable via signatures"
        ],
        5
    )
    
    # =========================================================================
    # SLIDE 6: INTRODUCTION - The AI Challenge
    # =========================================================================
    add_two_column_slide(
        prs,
        "INTRODUCTION - The AI Challenge (2020s)",
        "AI-Generated Phishing",
        [
            "Large Language Models (ChatGPT, Claude)",
            "Personalized, contextually sophisticated",
            "Grammatically flawless content",
            "Rapid evolution and adaptation"
        ],
        "Why Traditional Detection Fails",
        [
            "No technical artifacts or signatures",
            "Sophisticated linguistic patterns",
            "Contextually relevant content",
            "Continuous AI-driven evolution"
        ],
        6
    )
    
    # =========================================================================
    # SLIDE 7: PROBLEM STATEMENT - Core Problem
    # =========================================================================
    add_content_slide(
        prs,
        "PROBLEM STATEMENT - Core Problem",
        [
            "Develop intelligent, production-ready phishing detection system",
            "Real-time detection capability required",
            "Specific capability to identify AI-generated phishing attacks",
            "Must address reactive nature of traditional blacklists",
            "Need to handle sophisticated, evolving attack patterns",
            "Balance accuracy with resource efficiency"
        ],
        7
    )
    
    # =========================================================================
    # SLIDE 8: PROBLEM STATEMENT - Specific Challenges
    # =========================================================================
    add_table_slide(
        prs,
        "PROBLEM STATEMENT - Specific Challenges",
        [
            ["Challenge", "Impact", "Our Solution"],
            ["IDN/Homograph Attacks", "Cyrillic lookalike domains", "Punycode + mixed script analysis"],
            ["Typosquatting", "Brand impersonation", "50+ brand protection"],
            ["AI-Generated Content", "Evades traditional ML", "MLLM-based analysis"],
            ["Resource Efficiency", "High computational cost", "Tiered pipeline + cache"],
            ["Real-time Processing", "Millisecond latency needed", "4-tier architecture"]
        ],
        8
    )
    
    # =========================================================================
    # SLIDE 9: PROBLEM STATEMENT - Societal Impact
    # =========================================================================
    add_two_column_slide(
        prs,
        "PROBLEM STATEMENT - Societal Impact",
        "Financial & Personal Impact",
        [
            "Individual: Financial fraud, identity theft",
            "Businesses: Data breaches, reputational damage",
            "Economy: Billions in annual losses",
            "Trust: Erosion of digital infrastructure"
        ],
        "Research Gap",
        [
            "Existing systems: 85-92% accuracy",
            "No distinction AI vs traditional phishing",
            "Limited enterprise deployment options",
            "Need production-ready solutions"
        ],
        9
    )
    
    # =========================================================================
    # SLIDE 10: OBJECTIVES - Primary Objective
    # =========================================================================
    add_content_slide(
        prs,
        "OBJECTIVES - Primary Objective",
        [
            "Design and implement production-grade AI-powered phishing detection",
            "Use Multimodal Large Language Models",
            "Achieve >96% F1 score target",
            "Detect 4 distinct categories of phishing threats",
            "Real-time processing capability",
            "Enterprise-grade security features"
        ],
        10
    )
    
    # =========================================================================
    # SLIDE 11: OBJECTIVES - Technical Objectives
    # =========================================================================
    add_two_column_slide(
        prs,
        "OBJECTIVES - Technical Objectives",
        "Feature Engineering & ML",
        [
            "Extract 93 discriminative features",
            "Implement IDN/homograph detection",
            "Typosquatting protection for 50+ brands",
            "Train Random Forest classifier",
            "Achieve 99.8% F1 score",
            "MLflow for experiment tracking"
        ],
        "System & Security",
        [
            "4-tier detection architecture",
            "JWT-based API authentication",
            "Fernet encryption for sensitive data",
            "SSRF protection and input validation",
            "Rate limiting (100 req/min)",
            "REST API and browser extension"
        ],
        11
    )
    
    # =========================================================================
    # SLIDE 12: OBJECTIVES - Performance Targets
    # =========================================================================
    add_table_slide(
        prs,
        "OBJECTIVES - Performance Targets",
        [
            ["Metric", "Target", "Achieved"],
            ["F1 Score", ">96%", "99.8%"],
            ["Accuracy", ">96%", "99.6%"],
            ["False Positive Rate", "<1%", "<0.5%"],
            ["Detection Latency", "<2 seconds", "<2 seconds"],
            ["Throughput", "100+ URLs/minute", "100+ URLs/minute"]
        ],
        12
    )
    
    # =========================================================================
    # SLIDE 13: METHODOLOGY - Overview
    # =========================================================================
    add_content_slide(
        prs,
        "METHODOLOGY - Overall Approach",
        [
            "Systematic, data-driven approach combining ML and LLM",
            "Data Collection: 46,000+ URLs from PhishTank, OpenPhish",
            "Feature Extraction: 93 discriminative features",
            "Tiered Detection Pipeline for efficiency",
            "4-Category Classification System",
            "Production-ready deployment options"
        ],
        13
    )
    
    # =========================================================================
    # SLIDE 14: METHODOLOGY - Data Collection
    # =========================================================================
    add_table_slide(
        prs,
        "METHODOLOGY - Data Collection",
        [
            ["Source", "Type", "Count", "Purpose"],
            ["PhishTank", "Phishing URLs", "46,317", "Primary dataset"],
            ["OpenPhish", "Phishing URLs", "300", "Additional samples"],
            ["Custom Collection", "Legitimate URLs", "~10,000", "Baseline samples"],
            ["Total Dataset", "Combined", "46,617+", "Training & Testing"]
        ],
        14
    )
    
    # =========================================================================
    # SLIDE 15: METHODOLOGY - Feature Categories
    # =========================================================================
    add_table_slide(
        prs,
        "METHODOLOGY - Feature Categories (93 Total)",
        [
            ["Category", "Features", "Examples"],
            ["IDN/Punycode", "11", "xn-- detection, mixed scripts"],
            ["Host Analysis", "10", "Subdomain depth, suspicious TLDs"],
            ["URL Patterns", "28", "Char entropy, path analysis"],
            ["Security", "6", "SSRF, dangerous characters"],
            ["TLS/SSL", "11", "HTTPS, certificates, HSTS"],
            ["Composite", "3", "Risk scores, brand similarity"]
        ],
        15
    )
    
    # =========================================================================
    # SLIDE 16: METHODOLOGY - Typosquatting Detection
    # =========================================================================
    add_content_slide(
        prs,
        "METHODOLOGY - Typosquatting Detection",
        [
            "50+ major brands protected (Global + Indian)",
            "Levenshtein distance for edit distance calculation",
            "Character substitution pattern detection (0→o, 1→l)",
            "Homoglyph detection (Cyrillic lookalikes)",
            "TLD typosquatting (.pom→.com)",
            "Real-time brand protection database"
        ],
        16
    )
    
    # =========================================================================
    # SLIDE 17: METHODOLOGY - Machine Learning Pipeline
    # =========================================================================
    add_two_column_slide(
        prs,
        "METHODOLOGY - ML Classification",
        "Model: Random Forest",
        [
            "Excellent performance on tabular data",
            "Robust to overfitting",
            "Fast inference (<10ms)",
            "Interpretable feature importance",
            "200 estimators, max_depth=20"
        ],
        "Training Process",
        [
            "Data preprocessing & scaling",
            "5-fold stratified cross-validation",
            "Grid search optimization",
            "Full training set training",
            "Best model by F1 score"
        ],
        17
    )
    
    # =========================================================================
    # SLIDE 18: METHODOLOGY - MLLM Integration
    # =========================================================================
    add_two_column_slide(
        prs,
        "METHODOLOGY - MLLM Integration",
        "Model: Qwen2.5-3B-Instruct",
        [
            "4-bit quantization (4GB VRAM)",
            "Excellent instruction following",
            "Fast inference on consumer hardware",
            "Strong multimodal understanding",
            "Handles ambiguous cases"
        ],
        "Feature Transformation",
        [
            "URL structure analysis",
            "Screenshot description",
            "HTML content summary",
            "Security features checklist",
            "Unified text representation"
        ],
        18
    )
    
    # =========================================================================
    # SLIDE 19: METHODOLOGY - Tiered Pipeline Architecture
    # =========================================================================
    add_content_slide(
        prs,
        "METHODOLOGY - Tiered Detection Pipeline",
        [
            "TIER 1: Typosquatting Detection (<1ms) - Brand database + Levenshtein",
            "TIER 2: ML Classifier (<5ms) - 93 features → Random Forest",
            "TIER 3: MLLM Analysis (1-2s) - Multimodal feature transformation",
            "TIER 4: Result Cache (Redis) - 24-hour TTL for same URLs",
            "99% of traffic handled by Tier 2",
            "Only 1-5% requires expensive Tier 3 MLLM analysis"
        ],
        19
    )
    
    # =========================================================================
    # SLIDE 20: TECHNOLOGIES USED - Core Technologies
    # =========================================================================
    add_table_slide(
        prs,
        "TECHNOLOGIES USED - Core Stack",
        [
            ["Technology", "Version", "Purpose"],
            ["Python", "3.10+", "Primary language"],
            ["scikit-learn", "1.3+", "ML classification"],
            ["Qwen2.5-3B-Instruct", "Latest", "MLLM backbone"],
            ["Transformers", "4.35+", "Hugging Face library"],
            ["PyTorch", "2.0+", "Deep learning"],
            ["MLflow", "2.9+", "Experiment tracking"]
        ],
        20
    )
    
    # =========================================================================
    # SLIDE 21: TECHNOLOGIES USED - Backend & Security
    # =========================================================================
    add_two_column_slide(
        prs,
        "TECHNOLOGIES USED - Backend & Security",
        "Backend Technologies",
        [
            "FastAPI 0.109+ - REST API framework",
            "Uvicorn - ASGI server",
            "Pydantic - Data validation",
            "Redis 7.0+ - Caching layer",
            "Playwright 1.40+ - Web scraping"
        ],
        "Security Hardening",
        [
            "JWT authentication (24-hour tokens)",
            "Fernet AES-128 encryption",
            "SSRF protection (IP blocking)",
            "Rate limiting (100 req/min)",
            "TLS 1.3 enforcement"
        ],
        21
    )
    
    # =========================================================================
    # SLIDE 22: TECHNOLOGIES USED - Hardware & Development
    # =========================================================================
    add_table_slide(
        prs,
        "TECHNOLOGIES USED - Hardware & DevOps",
        [
            ["Component", "Minimum", "Recommended"],
            ["GPU", "NVIDIA RTX 3050 (4GB)", "NVIDIA RTX 4060 (8GB)"],
            ["RAM", "16GB", "32GB"],
            ["Storage", "50GB SSD", "100GB NVMe"],
            ["CPU", "4 cores", "8 cores"]
        ],
        22
    )
    
    # =========================================================================
    # SLIDE 23: RESULTS - Performance Metrics
    # =========================================================================
    add_table_slide(
        prs,
        "RESULTS - Classification Performance",
        [
            ["Metric", "Target", "Achieved", "Status"],
            ["F1 Score", ">96%", "99.8%", "✓ Exceeded"],
            ["Accuracy", ">96%", "99.6%", "✓ Exceeded"],
            ["Precision", ">95%", "99.7%", "✓ Exceeded"],
            ["Recall", ">95%", "99.8%", "✓ Exceeded"],
            ["False Positive Rate", "<1%", "<0.5%", "✓ Exceeded"]
        ],
        23
    )
    
    # =========================================================================
    # SLIDE 24: RESULTS - Detection by Category
    # =========================================================================
    add_table_slide(
        prs,
        "RESULTS - Performance by Category",
        [
            ["Category", "Precision", "Recall", "F1 Score"],
            ["LEGITIMATE", "99.5%", "99.8%", "99.7%"],
            ["PHISHING", "99.8%", "99.9%", "99.9%"],
            ["AI_GENERATED_PHISHING", "99.6%", "99.5%", "99.6%"],
            ["PHISHING_KIT", "99.7%", "99.6%", "99.7%"]
        ],
        24
    )
    
    # =========================================================================
    # SLIDE 25: RESULTS - Detection Capabilities
    # =========================================================================
    add_table_slide(
        prs,
        "RESULTS - Detection Capabilities",
        [
            ["Attack Type", "Detection Rate", "Examples"],
            ["Typosquatting", "99.9%", "paypa1.com, g00gle.com"],
            ["Punycode/IDN", "99.8%", "xn--paya1.com"],
            ["Homograph Attacks", "99.7%", "Cyrillic spoofing"],
            ["Phishing Kits", "99.6%", "Gophish, HiddenEye"],
            ["AI-Generated", "99.5%", "ChatGPT-crafted"],
            ["Brand Impersonation", "99.8%", "Fake PayPal, HDFC"]
        ],
        25
    )
    
    # =========================================================================
    # SLIDE 26: CONCLUSION - Key Accomplishments
    # =========================================================================
    add_content_slide(
        prs,
        "CONCLUSION - Key Accomplishments",
        [
            "Achieved 99.8% F1 score (3.8% above target)",
            "Implemented 93 discriminative features (+365% improvement)",
            "First production system for AI-generated phishing detection",
            "Novel IDN/homograph attack detection methodology",
            "Enterprise-grade security hardening",
            "Multiple deployment options (API, extension, desktop)"
        ],
        26
    )
    
    # =========================================================================
    # SLIDE 27: CONCLUSION - Project Statistics
    # =========================================================================
    add_table_slide(
        prs,
        "CONCLUSION - Project Statistics",
        [
            ["Metric", "Value"],
            ["Total Git Commits", "14+"],
            ["Files Created", "25+"],
            ["Lines of Code", "~6,000+"],
            ["Documentation", "2,500+ lines"],
            ["Test Classes", "14 suites"],
            ["Brand Protection", "50+ brands"],
            ["CVEs Patched", "8"]
        ],
        27
    )
    
    # =========================================================================
    # SLIDE 28: FUTURE SCOPE - Phase 2 Enhancements
    # =========================================================================
    add_content_slide(
        prs,
        "FUTURE SCOPE - Phase 2 Enhancements",
        [
            "Mobile Application (React Native) - iOS & Android",
            "Enhanced Browser Support - Firefox, Safari extensions",
            "Email Integration - Thunderbird, Gmail, Outlook plugins",
            "Threat Intelligence - Real-time feed integration",
            "Federated Learning - Privacy-preserving model updates",
            "Multi-language Phishing Detection - Global coverage"
        ],
        28
    )
    
    # =========================================================================
    # SLIDE 29: FUTURE SCOPE - Long-Term Vision
    # =========================================================================
    add_two_column_slide(
        prs,
        "FUTURE SCOPE - Long-Term Vision",
        "Advanced Capabilities",
        [
            "Fine-tuned domain-specific LLM",
            "Zero-shot attack detection",
            "Behavioral analysis",
            "Automated incident response"
        ],
        "Research Contributions",
        [
            "Academic publication",
            "Open-source core components",
            "Industry partnerships",
            "Collaborative research"
        ],
        29
    )
    
    # =========================================================================
    # SLIDE 30: THANK YOU
    # =========================================================================
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = PRIMARY_COLOR
    shape.line.fill.background()
    
    # Thank you text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "THANK YOU"
    p.font.size = Pt(64)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(12.333), Inches(1.5))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Questions & Discussion\n\nPhishing Guard v2.0 - AI-Powered Multimodal Phishing Detection System"
    p.font.size = Pt(24)
    p.font.color.rgb = ACCENT_COLOR
    p.alignment = PP_ALIGN.CENTER
    
    # Save presentation
    output_file = "Phishing_Guard_Presentation.pptx"
    prs.save(output_file)
    print(f"✓ Presentation saved as: {output_file}")
    print(f"✓ Total slides: {len(prs.slides)}")
    
    return output_file

if __name__ == "__main__":
    # Install required package if not installed
    try:
        from pptx import Presentation
    except ImportError:
        print("Installing python-pptx...")
        import subprocess
        subprocess.run(["pip", "install", "python-pptx"], check=True)
    
    create_presentation()
    print("\n✓ PowerPoint presentation created successfully!")
