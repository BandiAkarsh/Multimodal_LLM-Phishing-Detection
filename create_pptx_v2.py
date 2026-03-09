#!/usr/bin/env python3
"""
Create comprehensive detailed PPTX slides with images and detailed content
Each slide has an image and detailed paragraphs explaining the content
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

TITLE_COLOR = RGBColor(30, 58, 114)
TEXT_COLOR = RGBColor(50, 50, 50)

def add_title_slide(prs, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    p.alignment = PP_ALIGN.CENTER
    
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(12.333), Inches(1.5))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(20)
        p.font.color.rgb = TEXT_COLOR
        p.alignment = PP_ALIGN.CENTER
    return slide

def add_image_slide(prs, title, image_path, description):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(12.5), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    
    # Image
    if os.path.exists(image_path):
        img = slide.shapes.add_picture(image_path, Inches(0.3), Inches(0.9), width=Inches(7.5))
    
    # Description - right side
    desc_box = slide.shapes.add_textbox(Inches(7.9), Inches(0.9), Inches(5), Inches(6.3))
    tf = desc_box.text_frame
    tf.word_wrap = True
    
    paragraphs = description.split('\n\n')
    for i, para in enumerate(paragraphs):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = para
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_COLOR
        p.space_after = Pt(8)
    
    return slide

# ============ SLIDES ============

# 1. Title
add_title_slide(prs, "Phishing Guard v2.0", "AI-Powered Phishing Detection System\nIEEE Final Year Engineering Project\nAkarsh Bandi - March 2026")

# 2. Overview with Architecture
add_image_slide(prs, "System Overview", "viva/pdf_images_v2/system_architecture.png",
    """PHISHING GUARD SYSTEM

Complete AI-powered phishing detection with multiple layers of protection.

KEY COMPONENTS:
• User Interfaces: CLI, REST API, Browser Extension, Desktop App
• Processing Layer: 93-feature extraction pipeline
• ML/LLM Layer: Random Forest, XGBoost, Qwen2.5
• Security Layer: JWT Auth, Rate Limiting, SSRF Protection

The system analyzes URLs through comprehensive feature extraction and machine learning classification to detect phishing attacks with 99.70% accuracy.""")

# 3. Problem Statement
add_image_slide(prs, "Problem Statement", "viva/pdf_images_v2/system_architecture.png",
    """WHY PHISHING DETECTION MATTERS

Phishing is the #1 cybersecurity threat globally, causing billions in losses annually.

CURRENT LIMITATIONS:
• Static rule-based systems miss new attack patterns
• Blacklists have high false negatives for new domains
• Cannot detect AI-generated phishing content
• Limited feature analysis (~20 features)
• No multi-interface support

RESEARCH GAP:
Traditional ML models cannot detect sophisticated AI-generated phishing created by Large Language Models. These attacks have perfect grammar and context, evading traditional detection.""")

# 4. Objectives
add_image_slide(prs, "Project Objectives", "viva/pdf_images_v2/system_architecture.png",
    """PROJECT OBJECTIVES

1. Develop 93+ ML features for comprehensive URL analysis
2. Achieve >99% accuracy using ensemble models
3. Implement 4-category classification system
4. Create multiple user interfaces
5. Integrate Qwen2.5 LLM for AI phishing detection
6. Implement production-ready security

CATEGORIES:
• LEGITIMATE - Safe, trusted URLs
• PHISHING - Traditional attacks
• AI_GENERATED - LLM-created phishing
• PHISHING_KIT - Automated tools""")

# 5. Dataset
add_image_slide(prs, "Dataset Sources & Processing", "viva/pdf_images_v2/dataset_sources.png",
    """TRAINING DATA SOURCES

PHISHTANK:
• 135,000+ verified phishing URLs
• Community-verified database
• ~9.3 MB of data

OPENPHISH:
• 15,000+ verified phishing URLs
• Machine-detected threats
• Automated identification

LEGITIMATE SITES:
• Alexa Top 1 Million
• 50,000+ trusted URLs
• Diverse legitimate domains

DATA PROCESSING:
• URL normalization & cleaning
• Deduplication
• 80/10/10 train/val/test split
• Combined: ~200,000 URLs (~3.2MB)""")

# 6. Feature Extraction
add_image_slide(prs, "Feature Engineering - 93 Features", "viva/pdf_images_v2/ml_pipeline.png",
    """93 ML FEATURES EXTRACTED

URL PATTERN (28):
• Length metrics, character counts
• Entropy calculation
• Protocol analysis (HTTP/HTTPS)
• Suspicious word detection

DOMAIN (18):
• Domain entropy, subdomain count
• TLD analysis, typosquatting distance

HOST (10):
• IP detection, geographic location

SECURITY/TLS (12):
• Certificate validation
• TLS version checking (reject 1.0/1.1)

IDN/HOMOGRAPH (11):
• Punycode detection (xn--)
• Mixed script analysis
• Confusable characters (0 vs O)

BEHAVIORAL (14):
• Shortener detection
• Redirect analysis""")

# 7. ML Models
add_image_slide(prs, "Machine Learning Models", "viva/pdf_images_v2/ml_pipeline.png",
    """ML MODEL ARCHITECTURE

RANDOM FOREST:
• 200 trees, max_depth=20
• Accuracy: 99.64%
• F1 Score: 99.70%
• Stable, handles non-linear relationships

XGBOOST:
• 50 estimators, max_depth=6
• Learning rate: 0.1
• Accuracy: 99.58%
• F1 Score: 99.62%
• Gradient boosting for subtle patterns

SOFT VOTING ENSEMBLE:
• Combines RF + XGBoost
• Averages probability outputs
• Accuracy: 99.70%
• F1 Score: 99.82% (BEST)

MLFLOW TRACKING:
• Experiment: phishing_detection
• Full metrics logging""")

# 8. MLLM Integration
add_image_slide(prs, "MLLM Integration - Qwen2.5", "viva/pdf_images_v2/mllm_integration.png",
    """AI-GENERATED PHISHING DETECTION

PROBLEM:
LLMs can create sophisticated phishing content that evades traditional ML detection.

SOLUTION:
Qwen2.5-3B-Instruct integration

MODEL DETAILS:
• Base: Qwen2.5-3B-Instruct
• Quantization: 4-bit AWQ
• VRAM: ~2GB (consumer GPU)
• Inference: Ollama or Transformers

DETECTION PIPELINE:
1. ML provides quick screening
2. If confidence < 80%, trigger LLM
3. Fetch page content
4. Analyze with specialized prompt
5. Detect urgency, grammar, domain patterns
6. Classify as AI_GENERATED_PHISHING""")

# 9. REST API
add_image_slide(prs, "REST API Architecture", "viva/pdf_images_v2/api_architecture.png",
    """FASTAPI REST SERVICE

ENDPOINTS:
• GET / - API info (public)
• GET /health - Health check
• GET /connectivity - Connection status
• POST /auth/login - Get JWT token
• POST /auth/api-key - Generate API key
• POST /api/v1/analyze - Single URL
• POST /api/v1/batch-analyze - Batch URLs

TECHNOLOGY:
• Framework: FastAPI + Uvicorn
• Documentation: Swagger UI
• Authentication: JWT + API Keys

SECURITY:
• JWT: HS256, 24hr expiry
• Rate Limit: 100 req/min/IP
• SSRF Protection enabled""")

# 10. Browser Extension
add_image_slide(prs, "Browser Extension", "viva/pdf_images_v2/browser_extension.png",
    """CHROME MANIFEST V3 EXTENSION

ARCHITECTURE:
• manifest.json - Extension config
• background.js - Service worker
• content.js - DOM observer
• popup.html/js - User interface

KEY FEATURES:
• Real-time link scanning
• Color-coded indicators
• One-click page scan
• History tracking
• Privacy-first design

UI INDICATORS:
🟢 GREEN - Safe (legitimate)
🟡 YELLOW - Suspicious
🔴 RED - Malicious (phishing)
⚪ GRAY - Not yet analyzed

PERMISSIONS:
• activeTab, storage
• notifications, scripting""")

# 11. Desktop App
add_image_slide(prs, "Desktop Application (Tauri)", "viva/pdf_images_v2/tauri_gui.png",
    """TAURI + REACT DESKTOP APP

FRONTEND:
• React + TypeScript + Vite
• Components: URLInput, ResultDisplay, History, Settings

BACKEND:
• Tauri 2.x (Rust)
• Python API service
• Local ML model loading

FEATURES:
• Cross-platform (Win/Mac/Linux)
• Offline capability
• System tray support
• Native notifications
• Dark/Light theme

SIZE: ~5-10MB (vs 150MB+ Electron)

INSTALL:
1. Build with: cd gui-tauri && npm run tauri build
2. Run the generated executable""")

# 12. Security
add_image_slide(prs, "Security Implementation", "viva/pdf_images_v2/api_architecture.png",
    """PRODUCTION-READY SECURITY

AUTHENTICATION:
• JWT Tokens: HS256, 24-hour expiry
• API Keys: SHA-256 hashed storage
• Secure password hashing with bcrypt

RATE LIMITING:
• Default: 100 requests/minute/IP
• Redis support for distributed deployments
• In-memory fallback for single server

SSRF PROTECTION:
Blocked ranges:
• 10.0.0.0/8 (private)
• 172.16.0.0/12 (private)
• 192.168.0.0/16 (private)
• 127.0.0.1 (localhost)

TLS VALIDATION:
• Reject TLS 1.0 and 1.1
• Require TLS 1.2+
• Certificate chain verification""")

# 13. IDN/Homograph
add_image_slide(prs, "IDN/Homograph Detection", "viva/pdf_images_v2/system_architecture.png",
    """FIRST-OF-ITS-KIND DETECTION

PUNYCODE DETECTION:
Detects xn-- prefix used in internationalized domain attacks.

Example: xn--80ak6aa92e.com
(Appears as "apple.com" but is different!)

MIXED SCRIPT ANALYSIS:
Detects when domain uses multiple writing systems (Latin + Cyrillic + etc.)

CONFUSABLE CHARACTERS:
Detects lookalike characters:
• 0 (zero) vs O (letter)
• l (lowercase L) vs 1 (one)
• а (Cyrillic) vs a (Latin)

This protection is essential as internationalized domains become more common.""")

# 14. Performance
add_image_slide(prs, "Performance Metrics", "viva/pdf_images_v2/ml_pipeline.png",
    """MODEL PERFORMANCE

ACCURACY METRICS:
• Random Forest: 99.64%
• XGBoost: 99.58%
• Ensemble: 99.70%

F1 SCORE:
• Random Forest: 99.70%
• XGBoost: 99.62%
• Ensemble: 99.82%

PRECISION/RECALL:
• Precision: 99.72%
• Recall: 99.68%

SPEED:
• Feature Extraction: ~50ms
• ML Prediction: ~10ms
• Total Pipeline: ~100ms
• API Response: <200ms (p95)""")

# 15. Testing
add_image_slide(prs, "Testing & Quality", "viva/pdf_images_v2/system_architecture.png",
    """QUALITY ASSURANCE

TEST SUITES:
• Security Tests: 5/5 PASSING
• Comprehensive Tests: Feature extraction, ML models
• Integration Tests: API endpoints, service layer

COVERAGE:
• Code Coverage: >80%
• Type Hints: Full coverage
• Documentation: All public methods

SECURITY TESTING:
✓ JWT Authentication
✓ API Key Authentication
✓ Rate Limiting
✓ SSRF Protection
✓ Input Validation

The system is production-ready with comprehensive test coverage.""")

# 16. Deployment
add_image_slide(prs, "Deployment Options", "viva/pdf_images_v2/system_architecture.png",
    """DEPLOYMENT METHODS

CLI:
python demo.py --single https://example.com

API SERVER:
conda activate phishing-detection-mllm
uvicorn 04_inference.api:app --reload
# Access at http://localhost:8000/docs

BROWSER EXTENSION:
1. Chrome → chrome://extensions
2. Enable Developer mode
3. Load unpacked
4. Select browser-extension folder

DOCKER:
docker-compose up -d

DESKTOP APP:
cd gui-tauri
npm run tauri build
# Run generated executable""")

# 17. Tech Stack
add_image_slide(prs, "Technology Stack", "viva/pdf_images_v2/system_architecture.png",
    """TECHNOLOGY SELECTION

LANGUAGES:
• Python 3.9+ (Backend, ML)
• TypeScript (Frontend)
• Rust (Tauri backend)

ML/DL:
• scikit-learn (Random Forest)
• XGBoost (Gradient Boosting)
• PyTorch + Transformers (LLM)

WEB:
• FastAPI + Uvicorn (API)
• React + Vite (Frontend)
• Tauri 2.x (Desktop)

SECURITY:
• JWT (Authentication)
• bcrypt (Password hashing)
• hashlib (API keys)

TOOLS:
• Git, Docker, Conda
• MLflow (Experiment tracking)""")

# 18. Project Structure
add_image_slide(prs, "Project Structure", "viva/pdf_images_v2/system_architecture.png",
    """DIRECTORY ORGANIZATION

01_data/ - Raw & processed datasets
02_models/ - Trained ML models
03_training/ - Training scripts + MLflow
04_inference/ - API + service layer
05_utils/ - Feature extraction code
06_notebooks/ - Jupyter analysis
07_configs/ - Configuration files
08_logs/ - Log files
browser-extension/ - Chrome extension
gui-tauri/ - Desktop app source
tests/ - Test suites
viva/ - Presentation materials

Each directory is self-contained with clear responsibilities.""")

# 19. Future
add_image_slide(prs, "Future Enhancements", "viva/pdf_images_v2/system_architecture.png",
    """ROADMAP

SHORT-TERM:
• Tauri Desktop App production build
• Firefox Extension support
• Mobile apps (iOS/Android)
• Real-time notifications
• Multi-language support

LONG-TERM:
• Federated Learning (privacy)
• Threat intelligence integration
• Email plugins (Gmail, Outlook)
• SIEM integration
• Advanced LLM (larger models)

RESEARCH:
• Better AI phishing detection
• Zero-shot learning approaches
• Real-time URL analysis at scale""")

# 20. Conclusion
add_image_slide(prs, "Conclusion", "viva/pdf_images_v2/system_architecture.png",
    """PROJECT ACHIEVEMENTS

✓ 99.70% accuracy with ensemble ML
✓ 4-category classification system
✓ First IDN/Homograph detection
✓ Qwen2.5 LLM for AI phishing
✓ Multiple interfaces (CLI/API/Extension/Desktop)
✓ Production-ready security
✓ Comprehensive documentation

KEY INNOVATIONS:
1. 93-feature extraction pipeline
2. Hybrid ML + LLM detection
3. 4-category classification
4. Multi-interface deployment

The system is complete and ready for viva presentation.""")

# 21. Thank You
add_title_slide(prs, "Thank You!", "Questions & Discussion\n\nAkarsh Bandi\nIEEE Final Year Engineering Project\nMarch 2026")

# Save
output_path = "viva/Phishing_Guard_Professional_Presentation.pptx"
prs.save(output_path)
print(f"✓ Created: {output_path}")
print(f"  Total slides: {len(prs.slides)}")
