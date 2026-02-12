#!/usr/bin/env python3
"""
Phishing Guard v2.0 - IEEE Presentation Generator
Creates a professional 35-slide PowerPoint presentation

Usage:
    python create_presentation.py
    
Output:
    Phishing_Guard_IEEE_Presentation.pptx

Requirements:
    pip install python-pptx Pillow
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Color definitions
DEEP_BLUE = RGBColor(30, 58, 138)      # #1e3a8a
BLUE = RGBColor(59, 130, 246)          # #3b82f6
LIGHT_BLUE = RGBColor(219, 234, 254)   # #dbeafe
GREEN = RGBColor(16, 185, 129)         # #10b981
LIGHT_GREEN = RGBColor(209, 250, 229)  # #d1fae5
ORANGE = RGBColor(245, 158, 11)        # #f59e0b
RED = RGBColor(239, 68, 68)            # #ef4444
LIGHT_RED = RGBColor(254, 226, 226)    # #fee2e2
GOLD = RGBColor(251, 191, 36)          # #fbbf24
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(0, 0, 0)
GRAY = RGBColor(107, 114, 128)         # #6b7280

def create_presentation():
    """Create the 35-slide IEEE presentation"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title
    add_title_slide(prs)
    
    # Slide 2: Agenda
    add_agenda_slide(prs)
    
    # Slide 3: Abstract
    add_abstract_slide(prs)
    
    # Slide 4: Motivation
    add_motivation_slide(prs)
    
    # Slide 5: Problem Statement
    add_problem_slide(prs)
    
    # Slide 6: Threat Landscape
    add_threat_slide(prs)
    
    # Slide 7: Existing Solutions
    add_existing_slide(prs)
    
    # Slide 8: Research Gap
    add_gap_slide(prs)
    
    # Slide 9: Objectives
    add_objectives_slide(prs)
    
    # Slide 10: Key Innovations
    add_innovations_slide(prs)
    
    # Slide 11: Architecture
    add_architecture_slide(prs)
    
    # Slide 12: 4-Tier Pipeline
    add_pipeline_slide(prs)
    
    # Slide 13: Features
    add_features_slide(prs)
    
    # Slide 14: IDN Detection
    add_idn_slide(prs)
    
    # Slide 15: ML Classification
    add_ml_slide(prs)
    
    # Slide 16: MLLM
    add_mllm_slide(prs)
    
    # Slide 17: Web Scraping
    add_scraping_slide(prs)
    
    # Slide 18: Security
    add_security_slide(prs)
    
    # Slide 19: Multi-Interface
    add_interfaces_slide(prs)
    
    # Slide 20: Browser Extension
    add_extension_slide(prs)
    
    # Slide 21: Tech Stack
    add_techstack_slide(prs)
    
    # Slide 22: ML Technologies
    add_mltech_slide(prs)
    
    # Slide 23: Security Tech
    add_sectech_slide(prs)
    
    # Slide 24: Performance
    add_performance_slide(prs)
    
    # Slide 25: Comparison
    add_comparison_slide(prs)
    
    # Slide 26: Security Audit
    add_audit_slide(prs)
    
    # Slide 27: Test Coverage
    add_coverage_slide(prs)
    
    # Slide 28: Real-World Testing
    add_testing_slide(prs)
    
    # Slide 29: Competitive Matrix
    add_matrix_slide(prs)
    
    # Slide 30: UVP
    add_uvp_slide(prs)
    
    # Slide 31: Achievements
    add_achievements_slide(prs)
    
    # Slide 32: Impact
    add_impact_slide(prs)
    
    # Slide 33: Limitations
    add_limitations_slide(prs)
    
    # Slide 34: Future Work
    add_future_slide(prs)
    
    # Slide 35: Thank You
    add_thankyou_slide(prs)
    
    # Save presentation
    output_file = "Phishing_Guard_IEEE_Presentation.pptx"
    prs.save(output_file)
    print(f"✅ Presentation created: {output_file}")
    print(f"📊 Total slides: 35")
    print(f"🎓 Ready for IEEE submission!")
    return output_file

def add_title_slide(prs):
    """Slide 1: Title Slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Background
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = DEEP_BLUE
    background.line.fill.background()
    
    # Shield icon placeholder (text)
    shield = slide.shapes.add_textbox(Inches(5.5), Inches(1), Inches(2), Inches(1))
    tf = shield.text_frame
    p = tf.paragraphs[0]
    p.text = "🛡️"
    p.font.size = Pt(72)
    p.alignment = PP_ALIGN.CENTER
    
    # Title
    title = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(11.333), Inches(1))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Phishing Guard v2.0"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle = slide.shapes.add_textbox(Inches(1), Inches(3.3), Inches(11.333), Inches(0.8))
    tf = subtitle.text_frame
    p = tf.paragraphs[0]
    p.text = "Multimodal AI-Based Phishing Detection System with Enhanced Security"
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(200, 200, 200)
    p.alignment = PP_ALIGN.CENTER
    
    # Project info
    info = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(11.333), Inches(1.5))
    tf = info.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Final Year IEEE Project | Production-Grade Security Solution"
    p.font.size = Pt(18)
    p.font.color.rgb = GOLD
    p.alignment = PP_ALIGN.CENTER
    
    p = tf.add_paragraph()
    p.text = "\nPresented by: [Your Name]"
    p.font.size = Pt(16)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    p = tf.add_paragraph()
    p.text = "Institution: [Your University]"
    p.font.size = Pt(16)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    p = tf.add_paragraph()
    p.text = "February 2025"
    p.font.size = Pt(16)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

def add_agenda_slide(prs):
    """Slide 2: Agenda"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Presentation Agenda"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = DEEP_BLUE
    
    # Agenda items
    agenda_items = [
        ("1️⃣", "Introduction", "Project overview, motivation & problem statement"),
        ("2️⃣", "Background", "Threat landscape & existing solution gaps"),
        ("3️⃣", "Objectives", "Goals & key innovations"),
        ("4️⃣", "Methodology", "Architecture & 4-tier detection pipeline ⭐"),
        ("5️⃣", "Technologies", "Tech stack & tools"),
        ("6️⃣", "Results", "Performance metrics & validation ⭐"),
        ("7️⃣", "Comparison", "Competitive analysis"),
        ("8️⃣", "Conclusion", "Achievements & future work"),
    ]
    
    y_pos = 1.3
    for emoji, section, desc in agenda_items:
        # Number box
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y_pos), Inches(0.6), Inches(0.5))
        box.fill.solid()
        box.fill.fore_color.rgb = BLUE
        box.line.fill.background()
        
        num = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos + 0.05), Inches(0.6), Inches(0.4))
        tf = num.text_frame
        p = tf.paragraphs[0]
        p.text = emoji
        p.font.size = Pt(20)
        p.alignment = PP_ALIGN.CENTER
        
        # Section title
        sect = slide.shapes.add_textbox(Inches(1.3), Inches(y_pos), Inches(3), Inches(0.4))
        tf = sect.text_frame
        p = tf.paragraphs[0]
        p.text = section
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = DEEP_BLUE
        
        # Description
        desc_box = slide.shapes.add_textbox(Inches(4.5), Inches(y_pos), Inches(7), Inches(0.4))
        tf = desc_box.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(16)
        p.font.color.rgb = GRAY
        
        y_pos += 0.7
    
    # Highlights box
    highlight = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(6.5), Inches(12), Inches(0.7))
    highlight.fill.solid()
    highlight.fill.fore_color.rgb = LIGHT_BLUE
    highlight.line.color.rgb = BLUE
    
    hl_text = slide.shapes.add_textbox(Inches(0.5), Inches(6.6), Inches(12), Inches(0.5))
    tf = hl_text.text_frame
    p = tf.paragraphs[0]
    p.text = "⭐ Key Highlights: 99.8% F1 Score | 93 ML Features | 4-Category Classification"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = DEEP_BLUE
    p.alignment = PP_ALIGN.CENTER

def add_abstract_slide(prs):
    """Slide 3: Abstract"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Project Abstract"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = DEEP_BLUE
    
    # Problem box
    prob_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.2), Inches(6), Inches(2))
    prob_box.fill.solid()
    prob_box.fill.fore_color.rgb = LIGHT_RED
    prob_box.line.color.rgb = RED
    
    prob_title = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(5.5), Inches(0.4))
    tf = prob_title.text_frame
    p = tf.paragraphs[0]
    p.text = "🎯 Problem"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RED
    
    prob_text = slide.shapes.add_textbox(Inches(0.7), Inches(1.7), Inches(5.5), Inches(1.4))
    tf = prob_text.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "• 3.4B phishing emails/day\n• $4.88M average breach cost\n• AI-generated attacks rising\n• Visual spoofing undetected"
    p.font.size = Pt(14)
    p.font.color.rgb = BLACK
    
    # Solution box
    sol_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.2), Inches(6), Inches(2))
    sol_box.fill.solid()
    sol_box.fill.fore_color.rgb = LIGHT_GREEN
    sol_box.line.color.rgb = GREEN
    
    sol_title = slide.shapes.add_textbox(Inches(7), Inches(1.3), Inches(5.5), Inches(0.4))
    tf = sol_title.text_frame
    p = tf.paragraphs[0]
    p.text = "💡 Solution"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = GREEN
    
    sol_text = slide.shapes.add_textbox(Inches(7), Inches(1.7), Inches(5.5), Inches(1.4))
    tf = sol_text.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "• 4-tier multimodal pipeline\n• 99.8% F1 score achieved\n• IDN/Homograph detection\n• Production-grade security"
    p.font.size = Pt(14)
    p.font.color.rgb = BLACK
    
    # Metrics boxes
    metrics = [
        ("99.8%", "F1 Score", GREEN),
        ("93", "Features\n(+365%)", BLUE),
        ("4", "Category\nClassification", ORANGE),
        ("8", "CVEs\nPatched", RED),
    ]
    
    x_pos = 0.5
    for value, label, color in metrics:
        # Metric box
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x_pos), Inches(3.5), Inches(2.8), Inches(1.5))
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()
        
        # Value
        val_text = slide.shapes.add_textbox(Inches(x_pos), Inches(3.6), Inches(2.8), Inches(0.8))
        tf = val_text.text_frame
        p = tf.paragraphs[0]
        p.text = value
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        
        # Label
        lab_text = slide.shapes.add_textbox(Inches(x_pos), Inches(4.4), Inches(2.8), Inches(0.5))
        tf = lab_text.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        
        x_pos += 3.1
    
    # Innovation statement
    innov_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.3), Inches(12), Inches(1.5))
    innov_box.fill.solid()
    innov_box.fill.fore_color.rgb = GOLD
    innov_box.line.fill.background()
    
    innov_text = slide.shapes.add_textbox(Inches(0.7), Inches(5.4), Inches(11.6), Inches(1.3))
    tf = innov_text.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "🏆 INNOVATION"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = DEEP_BLUE
    
    p = tf.add_paragraph()
    p.text = "First system to combine IDN detection with AI-generated phishing classification in a production-ready platform"
    p.font.size = Pt(14)
    p.font.color.rgb = BLACK

def add_motivation_slide(prs):
    """Slide 4: Motivation"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "The Phishing Epidemic"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = DEEP_BLUE
    
    # Left side - Statistics
    stats_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.2), Inches(5.8), Inches(5.5))
    stats_box.fill.solid()
    stats_box.fill.fore_color.rgb = LIGHT_BLUE
    stats_box.line.color.rgb = BLUE
    
    stats_title = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(5.4), Inches(0.5))
    tf = stats_title.text_frame
    p = tf.paragraphs[0]
    p.text = "📊 By The Numbers"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = DEEP_BLUE
    
    stats = [
        ("📧", "3.4 Billion", "emails/day"),
        ("💰", "$4.88 Million", "avg cost/breach"),
        ("📈", "+1,265%", "AI phishing surge"),
        ("🎭", "82.6%", "contain AI content"),
        ("⚠️", "24.5%", "human detection rate"),
    ]
    
    y_pos = 2.0
    for emoji, value, desc in stats:
        # Emoji
        em = slide.shapes.add_textbox(Inches(0.8), Inches(y_pos), Inches(0.8), Inches(0.5))
        tf = em.text_frame
        p = tf.paragraphs[0]
        p.text = emoji
        p.font.size = Pt(24)
        
        # Value
        val = slide.shapes.add_textbox(Inches(1.5), Inches(y_pos), Inches(2.5), Inches(0.5))
        tf = val.text_frame
        p = tf.paragraphs[0]
        p.text = value
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RED
        
        # Description
        des = slide.shapes.add_textbox(Inches(4.0), Inches(y_pos), Inches(2), Inches(0.5))
        tf = des.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(14)
        p.font.color.rgb = GRAY
        
        y_pos += 0.9
    
    # Source
    src = slide.shapes.add_textbox(Inches(0.7), Inches(6.3), Inches(5.4), Inches(0.3))
    tf = src.text_frame
    p = tf.paragraphs[0]
    p.text = "Sources: APWG, IBM, KnowBe4 (2024-2025)"
    p.font.size = Pt(10)
    p.font.color.rgb = GRAY
    
    # Right side - Evolution
    evo_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.2), Inches(5.8), Inches(5.5))
    evo_box.fill.solid()
    evo_box.fill.fore_color.rgb = LIGHT_RED
    evo_box.line.color.rgb = RED
    
    evo_title = slide.shapes.add_textbox(Inches(7.0), Inches(1.3), Inches(5.4), Inches(0.5))
    tf = evo_title.text_frame
    p = tf.paragraphs[0]
    p.text = "⚠️ Attack Evolution"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RED
    
    evo_items = [
        ("2023: Basic Phishing", "✅ Detectable by traditional systems", GREEN),
        ("2024: AI-Generated", "⚠️ Perfect grammar, no typos\n❌ Traditional ML fails", ORANGE),
        ("2025: Visual Spoofing", "🔤 IDN Attacks (Cyrillic)\n❌ Invisible to users", RED),
    ]
    
    y_pos = 2.0
    for title, desc, color in evo_items:
        # Title
        t = slide.shapes.add_textbox(Inches(7.0), Inches(y_pos), Inches(5.4), Inches(0.4))
        tf = t.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = color
        
        # Description
        d = slide.shapes.add_textbox(Inches(7.2), Inches(y_pos + 0.4), Inches(5.2), Inches(0.8))
        tf = d.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(12)
        p.font.color.rgb = BLACK
        
        y_pos += 1.4
    
    # Call to action
    cta = slide.shapes.add_textbox(Inches(7.0), Inches(5.8), Inches(5.4), Inches(0.6))
    tf = cta.text_frame
    p = tf.paragraphs[0]
    p.text = "🆘 Next-Gen Detection NEEDED NOW!"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RED

def add_problem_slide(prs):
    """Slide 5: Problem Statement"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Critical Gaps in Current Solutions"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = DEEP_BLUE
    
    # Gap boxes
    gaps = [
        ("❌ GAP 1", "No IDN/\nHomograph\nProtection", RED),
        ("❌ GAP 2", "No AI-Generated\nDetection", ORANGE),
        ("❌ GAP 3", "Binary Only\nClassification", BLUE),
    ]
    
    x_pos = 0.5
    for gap_title, gap_desc, color in gaps:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x_pos), Inches(1.2), Inches(3.9), Inches(2.2))
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()
        
        # Title
        t = slide.shapes.add_textbox(Inches(x_pos), Inches(1.3), Inches(3.9), Inches(0.5))
        tf = t.text_frame
        p = tf.paragraphs[0]
        p.text = gap_title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        
        # Description
        d = slide.shapes.add_textbox(Inches(x_pos + 0.1), Inches(1.9), Inches(3.7), Inches(1.4))
        tf = d.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = gap_desc
        p.font.size = Pt(16)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        
        x_pos += 4.2
    
    # Gap 4 (bottom, full width)
    box4 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(3.7), Inches(12), Inches(1.5))
    box4.fill.solid()
    box4.fill.fore_color.rgb = RED
    box4.line.fill.background()
    
    t4 = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(12), Inches(0.4))
    tf = t4.text_frame
    p = tf.paragraphs[0]
    p.text = "❌ GAP 4: Inadequate Security"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    items = [
        "• No authentication",
        "• Plaintext credentials",
        "• No input validation",
        "• No rate limiting",
        "• No SSRF protection",
        "• Not production-ready",
    ]
    
    x_item = 0.8
    for item in items:
        it = slide.shapes.add_textbox(Inches(x_item), Inches(4.3), Inches(3.5), Inches(0.4))
        tf = it.text_frame
        p = tf.paragraphs[0]
        p.text = item
        p.font.size = Pt(12)
        p.font.color.rgb = WHITE
        x_item += 2.0
    
    # Research question
    rq_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.5), Inches(12), Inches(1.5))
    rq_box.fill.solid()
    rq_box.fill.fore_color.rgb = GOLD
    rq_box.line.fill.background()
    
    rq_title = slide.shapes.add_textbox(Inches(0.7), Inches(5.6), Inches(11.6), Inches(0.4))
    tf = rq_title.text_frame
    p = tf.paragraphs[0]
    p.text = "🔬 RESEARCH QUESTION"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = DEEP_BLUE
    
    rq_text = slide.shapes.add_textbox(Inches(0.7), Inches(6.0), Inches(11.6), Inches(0.8))
    tf = rq_text.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "How can we build a phishing detection system that addresses IDN attacks, AI-generated content, provides granular classification, AND meets enterprise security standards?"
    p.font.size = Pt(14)
    p.font.color.rgb = BLACK

# Continue with remaining slides... (simplified for brevity)
def add_threat_slide(prs):
    """Slide 6: Threat Landscape"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Modern Phishing Attack Vectors"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = DEEP_BLUE
    
    content = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5))
    tf = content.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "🎣 TRADITIONAL ATTACKS"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = BLUE
    
    p = tf.add_paragraph()
    p.text = "• Email spoofing • Typosquatting • Fake login pages • Urgency tactics"
    p.font.size = Pt(14)
    
    p = tf.add_paragraph()
    p.text = "\n🤖 AI-POWERED ATTACKS"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = ORANGE
    
    p = tf.add_paragraph()
    p.text = "• ChatGPT-generated content • Deepfake voice/video • Perfect grammar • Adaptive social engineering"
    p.font.size = Pt(14)
    
    p = tf.add_paragraph()
    p.text = "\n🌐 TECHNICAL ATTACKS"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RED
    
    p = tf.add_paragraph()
    p.text = "• IDN homographs • SSL/TLS abuse • Phishing kits (Gophish, Evilginx) • Man-in-the-middle"
    p.font.size = Pt(14)

def add_existing_slide(prs):
    """Slide 7: Existing Solutions"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Current Solutions & Limitations"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = DEEP_BLUE
    
    # Simple table representation
    table_data = [
        ["Feature", "PhishTank", "Google SB", "Traditional ML", "Phishing Guard"],
        ["Real-time Analysis", "❌", "❌", "✅", "✅"],
        ["AI Detection", "❌", "❌", "❌", "✅"],
        ["IDN Protection", "❌", "❌", "❌", "✅⭐"],
        ["4-Category Output", "❌", "❌", "❌", "✅⭐"],
        ["93 ML Features", "N/A", "N/A", "~20", "93⭐"],
        ["Enterprise Security", "N/A", "Enterprise", "Basic", "✅⭐"],
    ]
    
    y_pos = 1.5
    for row in table_data:
        x_pos = 0.5
        for cell in row:
            cell_box = slide.shapes.add_textbox(Inches(x_pos), Inches(y_pos), Inches(2.8), Inches(0.4))
            tf = cell_box.text_frame
            p = tf.paragraphs[0]
            p.text = cell
            p.font.size = Pt(12) if y_pos > 1.5 else Pt(14)
            p.font.bold = True if y_pos == 1.5 else False
            
            if "✅" in cell:
                p.font.color.rgb = GREEN
            elif "❌" in cell:
                p.font.color.rgb = RED
            elif "⭐" in cell:
                p.font.color.rgb = GOLD
            
            x_pos += 2.9
        y_pos += 0.5

def add_gap_slide(prs):
    """Slide 8: Research Gap"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "The Research Gap"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = DEEP_BLUE
    
    content = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5))
    tf = content.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "CURRENT SYSTEMS:"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RED
    
    p = tf.add_paragraph()
    p.text = "• 20 Features • Binary Classification (Safe/Unsafe) • No Security Hardening"
    p.font.size = Pt(14)
    
    p = tf.add_paragraph()
    p.text = "\nIDEAL SYSTEM:"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = GREEN
    
    p = tf.add_paragraph()
    p.text = "• 93 Features • 4-Category Classification • Enterprise Security"
    p.font.size = Pt(14)
    
    p = tf.add_paragraph()
    p.text = "\nSPECIFIC GAPS IDENTIFIED:"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = ORANGE
    
    gaps = [
        "1. Visual Spoofing: No Unicode/IDN detection",
        "2. AI Content: Cannot detect GPT-generated attacks",
        "3. Granularity: Binary only (no threat levels)",
        "4. Security: Not production-ready"
    ]
    
    for gap in gaps:
        p = tf.add_paragraph()
        p.text = gap
        p.font.size = Pt(14)
    
    p = tf.add_paragraph()
    p.text = "\n✅ OUR CONTRIBUTION: First unified system addressing ALL gaps!"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = GREEN

def add_objectives_slide(prs):
    """Slide 9: Objectives"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Project Objectives"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = DEEP_BLUE
    
    objectives = [
        ("🎯", "Build 4-Tier Detection Pipeline", "Static analysis, ML classification, MLLM reasoning, Web scraping"),
        ("🤖", "Detect AI-Generated Phishing", "MLLM integration, Linguistic analysis, Pattern recognition"),
        ("🔤", "Prevent IDN/Homograph Attacks", "Punycode detection, Mixed script analysis, Confusable mapping"),
        ("🔒", "Production-Grade Security", "JWT authentication, Rate limiting, SSRF protection"),
    ]
    
    y_pos = 1.3
    for emoji, obj_title, obj_desc in objectives:
        # Box
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y_pos), Inches(12), Inches(1.2))
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_BLUE
        box.line.color.rgb = BLUE
        
        # Emoji
        em = slide.shapes.add_textbox(Inches(0.7), Inches(y_pos + 0.1), Inches(0.8), Inches(0.6))
        tf = em.text_frame
        p = tf.paragraphs[0]
        p.text = emoji
        p.font.size = Pt(32)
        
        # Title
        t = slide.shapes.add_textbox(Inches(1.5), Inches(y_pos + 0.1), Inches(5), Inches(0.5))
        tf = t.text_frame
        p = tf.paragraphs[0]
        p.text = obj_title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = DEEP_BLUE
        
        # Description
        d = slide.shapes.add_textbox(Inches(1.5), Inches(y_pos + 0.55), Inches(10), Inches(0.6))
        tf = d.text_frame
        p = tf.paragraphs[0]
        p.text = obj_desc
        p.font.size = Pt(12)
        p.font.color.rgb = GRAY
        
        y_pos += 1.4

def add_innovations_slide(prs):
    """Slide 10: Key Innovations"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "5 Key Innovations (Unique Selling Points)"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = DEEP_BLUE
    
    innovations = [
        ("🥇", "IDN Detection", "FIRST academic implementation with full Unicode analysis"),
        ("🤖", "AI-Generated Classification", "Addresses ChatGPT-powered threat landscape"),
        ("📊", "365% Feature Improvement", "93 features vs industry standard 20"),
        ("🔒", "Enterprise Security", "8 CVEs patched, production-ready hardening"),
        ("🌐", "Multi-Modal Architecture", "4-tier ensemble with latency optimization"),
    ]
    
    y_pos = 1.2
    for emoji, inv_title, inv_desc in innovations:
        # Gold box
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y_pos), Inches(12), Inches(0.95))
        box.fill.solid()
        box.fill.fore_color.rgb = GOLD
        box.line.fill.background()
        
        # Content
        content = slide.shapes.add_textbox(Inches(0.7), Inches(y_pos + 0.1), Inches(11.6), Inches(0.75))
        tf = content.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = f"{emoji} {inv_title}"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = DEEP_BLUE
        
        p = tf.add_paragraph()
        p.text = inv_desc
        p.font.size = Pt(12)
        p.font.color.rgb = BLACK
        
        y_pos += 1.1

def add_architecture_slide(prs):
    """Slide 11: System Architecture"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "System Architecture - 4 Layers"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = DEEP_BLUE
    
    layers = [
        ("LAYER 1: USER INTERFACES", "CLI • Web API • Browser Extension • Desktop GUI • Email Scanner", LIGHT_BLUE, BLUE),
        ("LAYER 2: API & SECURITY", "FastAPI • JWT Auth • Rate Limiting • Input Validation • CORS", LIGHT_GREEN, GREEN),
        ("LAYER 3: DETECTION PIPELINE", "Tier 1: Static • Tier 2: ML • Tier 3: MLLM • Tier 4: Scraping", LIGHT_RED, RED),
        ("LAYER 4: DATA & MODELS", "Random Forest • MLflow Registry • Redis Cache • Feature Extractor", RGBColor(254, 243, 199), ORANGE),
    ]
    
    y_pos = 1.2
    for layer_title, layer_content, bg_color, border_color in layers:
        # Box
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y_pos), Inches(12), Inches(1.2))
        box.fill.solid()
        box.fill.fore_color.rgb = bg_color
        box.line.color.rgb = border_color
        box.line.width = Pt(2)
        
        # Title
        t = slide.shapes.add_textbox(Inches(0.7), Inches(y_pos + 0.1), Inches(11.6), Inches(0.4))
        tf = t.text_frame
        p = tf.paragraphs[0]
        p.text = layer_title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = border_color
        
        # Content
        c = slide.shapes.add_textbox(Inches(0.7), Inches(y_pos + 0.5), Inches(11.6), Inches(0.6))
        tf = c.text_frame
        p = tf.paragraphs[0]
        p.text = layer_content
        p.font.size = Pt(14)
        p.font.color.rgb = BLACK
        
        y_pos += 1.4

def add_pipeline_slide(prs):
    """Slide 12: 4-Tier Pipeline"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "4-Tier Multimodal Detection Pipeline"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = DEEP_BLUE
    
    tiers = [
        ("TIER 1", "STATIC", "~10ms", GREEN, "URL parsing • Typosquatting • IDN detection • Whitelist check"),
        ("TIER 2", "ML CLASSIFIER", "~100ms", BLUE, "93 features • Random Forest • Probability • Feature importance"),
        ("TIER 3", "MLLM ANALYSIS", "~2s", ORANGE, "Qwen2.5-3B • AI detection • Linguistic patterns • Optional tier"),
        ("TIER 4", "WEB SCRAPING", "Variable", RED, "Playwright • Screenshots • Form analysis • Toolkit fingerprint"),
    ]
    
    y_pos = 1.2
    for tier_num, tier_name, latency, color, features in tiers:
        # Box
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y_pos), Inches(12), Inches(1.1))
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()
        
        # Number and name
        t = slide.shapes.add_textbox(Inches(0.7), Inches(y_pos + 0.1), Inches(4), Inches(0.5))
        tf = t.text_frame
        p = tf.paragraphs[0]
        p.text = f"{tier_num}: {tier_name}"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = WHITE
        
        # Latency
        lat = slide.shapes.add_textbox(Inches(5), Inches(y_pos + 0.1), Inches(2), Inches(0.5))
        tf = lat.text_frame
        p = tf.paragraphs[0]
        p.text = latency
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = WHITE
        
        # Features
        feat = slide.shapes.add_textbox(Inches(0.7), Inches(y_pos + 0.55), Inches(11.6), Inches(0.5))
        tf = feat.text_frame
        p = tf.paragraphs[0]
        p.text = features
        p.font.size = Pt(11)
        p.font.color.rgb = WHITE
        
        y_pos += 1.25

def add_features_slide(prs):
    """Slide 13: Feature Extraction"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "93 ML Features - 365% Improvement"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = DEEP_BLUE
    
    # Feature categories
    categories = [
        ("IDN/Unicode", "11", RED),
        ("Host Analysis", "10", BLUE),
        ("URL Patterns", "28", GREEN),
        ("Security", "6", ORANGE),
        ("TLS/SSL", "11", RGBColor(139, 92, 246)),
        ("Composite", "3", RGBColor(236, 72, 153)),
        ("Other", "24", GRAY),
    ]
    
    y_pos = 1.3
    for cat_name, cat_count, color in categories:
        # Category box
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y_pos), Inches(5.5), Inches(0.6))
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()
        
        # Name
        name = slide.shapes.add_textbox(Inches(0.7), Inches(y_pos + 0.1), Inches(3.5), Inches(0.4))
        tf = name.text_frame
        p = tf.paragraphs[0]
        p.text = cat_name
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE
        
        # Count
        count = slide.shapes.add_textbox(Inches(4.5), Inches(y_pos + 0.1), Inches(1.3), Inches(0.4))
        tf = count.text_frame
        p = tf.paragraphs[0]
        p.text = cat_count
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.RIGHT
        
        y_pos += 0.7
    
    # Comparison box
    comp_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.5), Inches(1.3), Inches(6), Inches(2))
    comp_box.fill.solid()
    comp_box.fill.fore_color.rgb = LIGHT_BLUE
    comp_box.line.color.rgb = BLUE
    
    comp_text = slide.shapes.add_textbox(Inches(6.7), Inches(1.4), Inches(5.6), Inches(1.8))
    tf = comp_text.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "COMPARISON"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = DEEP_BLUE
    
    p = tf.add_paragraph()
    p.text = "\nTraditional: 20 features"
    p.font.size = Pt(14)
    p.font.color.rgb = GRAY
    
    p = tf.add_paragraph()
    p.text = "Phishing Guard: 93 features"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = GREEN
    
    p = tf.add_paragraph()
    p.text = "\n+365% Improvement!"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = GOLD
    
    # Key features box
    key_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.5), Inches(3.5), Inches(6), Inches(3.5))
    key_box.fill.solid()
    key_box.fill.fore_color.rgb = GOLD
    key_box.line.fill.background()
    
    key_text = slide.shapes.add_textbox(Inches(6.7), Inches(3.6), Inches(5.6), Inches(3.3))
    tf = key_text.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "KEY NOVEL FEATURES"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = DEEP_BLUE
    
    features = [
        "• has_punycode - IDN detection",
        "• mixed_scripts - Multi-script",
        "• confusable_count - Lookalikes",
        "• hsts_enabled - Security",
        "• ct_logs_found - Transparency",
        "• domain_age_days - Temporal",
    ]
    
    for feat in features:
        p = tf.add_paragraph()
        p.text = feat
        p.font.size = Pt(11)
        p.font.color.rgb = BLACK

def add_idn_slide(prs):
    """Slide 14: IDN Detection - NOVEL CONTRIBUTION"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "IDN/Homograph Attack Detection - NOVEL! ⭐"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RED
    
    # Problem box
    prob_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.2), Inches(12), Inches(1.5))
    prob_box.fill.solid()
    prob_box.fill.fore_color.rgb = LIGHT_RED
    prob_box.line.color.rgb = RED
    
    prob_text = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(11.6), Inches(1.3))
    tf = prob_text.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "THE PROBLEM: Visual Spoofing"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RED
    
    p = tf.add_paragraph()
    p.text = "Legitimate: paypal.com (ASCII)"
    p.font.size = Pt(14)
    p.font.color.rgb = GREEN
    
    p = tf.add_paragraph()
    p.text = "Malicious: раураl.com (Cyrillic 'р' and 'а') - LOOKS IDENTICAL!"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RED
    
    # Solution box
    sol_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(2.9), Inches(12), Inches(4))
    sol_box.fill.solid()
    sol_box.fill.fore_color.rgb = LIGHT_GREEN
    sol_box.line.color.rgb = GREEN
    
    sol_text = slide.shapes.add_textbox(Inches(0.7), Inches(3.0), Inches(11.6), Inches(3.8))
    tf = sol_text.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "OUR SOLUTION: 11 Specialized IDN Features"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = GREEN
    
    features = [
        "1. has_punycode - Detects xn-- prefix",
        "2. mixed_scripts - Latin + Cyrillic detection",
        "3. confusable_count - 25+ lookalike characters",
        "4. unicode_category - Letter/Mark/Number/Symbol",
        "5. script_blocks - Identifies all scripts used",
        "6. idn_risk_score - Composite threat score",
        "7. homograph_ratio - Suspicious character %",
        "8. punycode_ratio - Encoded character %",
        "9. decoded_domain - Original Unicode form",
        "10. visual_similarity - Confusion assessment",
        "11. brand_confusion - Impersonation detection",
    ]
    
    for feat in features:
        p = tf.add_paragraph()
        p.text = feat
        p.font.size = Pt(12)
        p.font.color.rgb = BLACK
    
    # Result
    result = slide.shapes.add_textbox(Inches(0.5), Inches(6.2), Inches(12), Inches(0.6))
    tf = result.text_frame
    p = tf.paragraphs[0]
    p.text = "✅ RESULT: Blocks visual spoofing attacks that bypass all other filters!"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = GREEN

def add_ml_slide(prs):
    """Slide 15: ML Classification"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Machine Learning Classification"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = DEEP_BLUE
    
    # Model box
    model_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.2), Inches(6), Inches(2.5))
    model_box.fill.solid()
    model_box.fill.fore_color.rgb = LIGHT_BLUE
    model_box.line.color.rgb = BLUE
    
    model_text = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(5.6), Inches(2.3))
    tf = model_text.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "MODEL ARCHITECTURE"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = DEEP_BLUE
    
    specs = [
        "• Algorithm: Random Forest",
        "• Estimators: 200",
        "• Features: 93 input features",
        "• Classes: 4 categories",
        "• Training Data: 10,000+ URLs",
    ]
    
    for spec in specs:
        p = tf.add_paragraph()
        p.text = spec
        p.font.size = Pt(12)
        p.font.color.rgb = BLACK
    
    # Metrics box
    met_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.2), Inches(5.5), Inches(2.5))
    met_box.fill.solid()
    met_box.fill.fore_color.rgb = GOLD
    met_box.line.fill.background()
    
    met_text = slide.shapes.add_textbox(Inches(7.0), Inches(1.3), Inches(5.1), Inches(2.3))
    tf = met_text.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "PERFORMANCE"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = DEEP_BLUE
    
    p = tf.add_paragraph()
    p.text = "\nF1 Score: 99.8%"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = GREEN
    
    p = tf.add_paragraph()
    p.text = "Accuracy: 99.6%"
    p.font.size = Pt(16)
    p.font.color.rgb = BLACK
    
    p = tf.add_paragraph()
    p.text = "Inference: ~100ms"
    p.font.size = Pt(16)
    p.font.color.rgb = BLACK
    
    # 4 categories box
    cat_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.0), Inches(12), Inches(2.8))
    cat_box.fill.solid()
    cat_box.fill.fore_color.rgb = LIGHT_GREEN
    cat_box.line.color.rgb = GREEN
    
    cat_text = slide.shapes.add_textbox(Inches(0.7), Inches(4.1), Inches(11.6), Inches(2.6))
    tf = cat_text.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "4-CATEGORY CLASSIFICATION (Industry-First!)"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = DEEP_BLUE
    
    cats = [
        "✅ Legitimate - Safe URLs, whitelisted domains",
        "🔴 Phishing - Traditional spoofing attacks",
        "🟠 AI-Generated - ChatGPT/AI-created content (NEW!)",
        "🚨 Phishing Kit - Toolkit-generated (Gophish, etc.)",
    ]
    
    for cat in cats:
        p = tf.add_paragraph()
        p.text = cat
        p.font.size = Pt(13)
        p.font.color.rgb = BLACK

# Continue with simplified versions of remaining slides...
def add_mllm_slide(prs):
    """Slide 16: MLLM Integration"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "MLLM Integration for AI-Generated Detection"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = DEEP_BLUE
    
    content = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5))
    tf = content.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Qwen2.5-3B-Instruct Model"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = ORANGE
    
    specs = [
        "• 3 billion parameters",
        "• 4-bit quantized (2GB memory)",
        "• ~2 seconds inference time",
        "• Optional Tier 3 analysis",
    ]
    
    for spec in specs:
        p = tf.add_paragraph()
        p.text = spec
        p.font.size = Pt(14)
    
    p = tf.add_paragraph()
    p.text = "\nWHY MLLM?"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = DEEP_BLUE
    
    p = tf.add_paragraph()
    p.text = "AI-generated phishing has perfect grammar, no typos. Traditional ML fails. MLLM detects linguistic AI patterns."
    p.font.size = Pt(14)
    
    p = tf.add_paragraph()
    p.text = "\n✅ Impact: +3.2% accuracy on AI-generated samples"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = GREEN

def add_scraping_slide(prs):
    """Slide 17: Web Scraping"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Web Scraping & Toolkit Fingerprinting"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = DEEP_BLUE
    
    content = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5))
    tf = content.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Playwright Headless Browser"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = BLUE
    
    features = [
        "• Dynamic content analysis",
        "• Screenshot capture for MLLM",
        "• Form field analysis",
        "• External link mapping",
    ]
    
    for feat in features:
        p = tf.add_paragraph()
        p.text = feat
        p.font.size = Pt(14)
    
    p = tf.add_paragraph()
    p.text = "\nTOOLKIT SIGNATURES DETECTED:"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RED
    
    tools = [
        "• Gophish - Track ID, rid parameter",
        "• Evilginx2 - Lure path, custom headers",
        "• HiddenEye - Template patterns",
        "• SocialFish - Redirect patterns",
    ]
    
    for tool in tools:
        p = tf.add_paragraph()
        p.text = tool
        p.font.size = Pt(14)

def add_security_slide(prs):
    """Slide 18: Security Architecture"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Enterprise Security Architecture"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = DEEP_BLUE
    
    content = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5))
    tf = content.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Defense in Depth - 8 Security Layers"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RED
    
    layers = [
        "1. JWT Authentication (24hr tokens)",
        "2. Rate Limiting (100 req/min)",
        "3. Input Validation (RFC 3986)",
        "4. SSRF Protection (Private IP blocking)",
        "5. TLS 1.3 Enforcement",
        "6. Credential Encryption (AES-128)",
        "7. CORS Protection",
        "8. Audit Logging",
    ]
    
    for layer in layers:
        p = tf.add_paragraph()
        p.text = layer
        p.font.size = Pt(14)
    
    p = tf.add_paragraph()
    p.text = "\n🏆 8 CVE-Level Vulnerabilities Patched!"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = GREEN

def add_interfaces_slide(prs):
    """Slide 19: Multi-Interface"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Multi-Interface Deployment"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = DEEP_BLUE
    
    interfaces = [
        ("1️⃣ CLI", "Color-coded output, batch processing, JSON export"),
        ("2️⃣ REST API", "FastAPI, JWT auth, Swagger docs"),
        ("3️⃣ Browser Extension", "Real-time scanning, visual highlighting"),
        ("4️⃣ Desktop GUI", "Tauri app, 3.8MB, cross-platform"),
        ("5️⃣ Email Scanner", "IMAP integration, attachment analysis"),
    ]
    
    y_pos = 1.3
    for iface, desc in interfaces:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y_pos), Inches(12), Inches(0.9))
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_BLUE
        box.line.color.rgb = BLUE
        
        text = slide.shapes.add_textbox(Inches(0.7), Inches(y_pos + 0.15), Inches(11.6), Inches(0.6))
        tf = text.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = f"{iface}: {desc}"
        p.font.size = Pt(14)
        p.font.color.rgb = BLACK
        
        y_pos += 1.05

def add_extension_slide(prs):
    """Slide 20: Browser Extension"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Browser Extension - Real-Time Protection"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = DEEP_BLUE
    
    content = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5))
    tf = content.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Visual Threat Highlighting:"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = DEEP_BLUE
    
    highlights = [
        "🟢 Green underline - Safe/Legitimate",
        "🔴 Red border - Phishing detected",
        "🟠 Orange dashed - AI-Generated content",
        "🔴 Dark red - Phishing Kit detected",
    ]
    
    for hl in highlights:
        p = tf.add_paragraph()
        p.text = hl
        p.font.size = Pt(14)
    
    p = tf.add_paragraph()
    p.text = "\nFeatures:"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = DEEP_BLUE
    
    feats = [
        "• Automatic scanning on page load",
        "• Dynamic content detection",
        "• Desktop notifications",
        "• Configurable API URL",
        "• Chrome/Brave/Edge support",
    ]
    
    for feat in feats:
        p = tf.add_paragraph()
        p.text = feat
        p.font.size = Pt(14)

def add_techstack_slide(prs):
    """Slide 21: Tech Stack"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Technology Stack Overview"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = DEEP_BLUE
    
    content = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5))
    tf = content.text_frame
    tf.word_wrap = True
    
    stacks = [
        ("🐍 Core", "Python 3.11, FastAPI, Pydantic, PyJWT"),
        ("🤖 ML/AI", "scikit-learn, PyTorch, Transformers, MLflow"),
        ("🌐 Web", "Playwright, BeautifulSoup, requests"),
        ("🗄️ Data", "pandas, numpy, Redis"),
        ("🛠️ DevOps", "Docker, GitHub Actions, pytest"),
        ("🖥️ Frontend", "Tauri, React, TypeScript, Rust"),
    ]
    
    y_pos = 1.5
    for stack_name, stack_tech in stacks:
        p = tf.paragraphs[0] if y_pos == 1.5 else tf.add_paragraph()
        p.text = f"{stack_name}: {stack_tech}"
        p.font.size = Pt(16)
        
        if y_pos == 1.5:
            y_pos += 0.6
        else:
            y_pos += 0.6
    
    p = tf.add_paragraph()
    p.text = "\nTotal: 30+ Technologies | 10,775 Lines of Code"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = GOLD

def add_mltech_slide(prs):
    """Slide 22: ML Technologies"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "ML/AI Technologies"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = DEEP_BLUE
    
    content = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5))
    tf = content.text_frame
    tf.word_wrap = True
    
    models = [
        ("Random Forest", "200 estimators, 93 features, 99.8% F1"),
        ("Qwen2.5-3B", "3B params, 4-bit quantized, AI detection"),
        ("MLflow", "Model versioning, experiment tracking"),
    ]
    
    for model, spec in models:
        p = tf.paragraphs[0] if model == "Random Forest" else tf.add_paragraph()
        p.text = f"\n{model}:"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = BLUE
        
        p = tf.add_paragraph()
        p.text = spec
        p.font.size = Pt(14)

def add_sectech_slide(prs):
    """Slide 23: Security Tech"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Security & Deployment Technologies"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = DEEP_BLUE
    
    content = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5))
    tf = content.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Security Stack:"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RED
    
    sec_items = [
        "• PyJWT, cryptography, keyring",
        "• Pydantic validation, RFC 3986",
        "• Redis rate limiting, TLS 1.3",
    ]
    
    for item in sec_items:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(14)
    
    p = tf.add_paragraph()
    p.text = "\nDeployment Options:"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = BLUE
    
    dep_items = [
        "• Docker + docker-compose",
        "• Kubernetes ready",
        "• Systemd service (166KB)",
    ]
    
    for item in dep_items:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(14)

def add_performance_slide(prs):
    """Slide 24: Performance - KEY SLIDE"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = LIGHT_GREEN
    bg.line.fill.background()
    
    # Title
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Performance Metrics"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = DEEP_BLUE
    
    # Big metric
    big = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(1.5))
    tf = big.text_frame
    p = tf.paragraphs[0]
    p.text = "99.8%"
    p.font.size = Pt(96)
    p.font.bold = True
    p.font.color.rgb = GREEN
    p.alignment = PP_ALIGN.CENTER
    
    sub = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(12), Inches(0.5))
    tf = sub.text_frame
    p = tf.paragraphs[0]
    p.text = "F1 SCORE"
    p.font.size = Pt(36)
    p.font.color.rgb = DEEP_BLUE
    p.alignment = PP_ALIGN.CENTER
    
    # Metrics boxes
    metrics = [
        ("Accuracy", "99.6%", BLUE),
        ("Precision", "99.7%", BLUE),
        ("Recall", "99.9%", BLUE),
    ]
    
    x_pos = 1.5
    for name, value, color in metrics:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x_pos), Inches(3.5), Inches(3.2), Inches(1.2))
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()
        
        v = slide.shapes.add_textbox(Inches(x_pos), Inches(3.6), Inches(3.2), Inches(0.7))
        tf = v.text_frame
        p = tf.paragraphs[0]
        p.text = value
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        
        n = slide.shapes.add_textbox(Inches(x_pos), Inches(4.3), Inches(3.2), Inches(0.4))
        tf = n.text_frame
        p = tf.paragraphs[0]
        p.text = name
        p.font.size = Pt(16)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        
        x_pos += 3.5
    
    # Bottom metrics
    bottom = slide.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(12), Inches(1.5))
    tf = bottom.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Additional Metrics:"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = DEEP_BLUE
    
    add_metrics = [
        "• False Positive Rate: <0.5% (Industry: 2-5%)",
        "• Latency: <150ms (75% faster)",
        "• Throughput: 100+ URLs/minute",
        "• Test Coverage: 100% (security-critical)",
    ]
    
    for metric in add_metrics:
        p = tf.add_paragraph()
        p.text = metric
        p.font.size = Pt(14)

def add_comparison_slide(prs):
    """Slide 25: Comparison"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Competitive Analysis"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = DEEP_BLUE
    
    content = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5))
    tf = content.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Multi-Dimensional Comparison (5 Dimensions)"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = DEEP_BLUE
    
    dims = [
        "Phishing Guard: 10/10 F1, 10/10 Features, 10/10 Security, 10/10 Granularity, 9/10 Latency",
        "PhishTank: Limited features, reactive only",
        "Google SB: Good coverage, limited features",
        "Traditional ML: Basic detection, poor granularity",
    ]
    
    for dim in dims:
        p = tf.add_paragraph()
        p.text = dim
        p.font.size = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "\n✅ Phishing Guard leads in 4 out of 5 dimensions!"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = GREEN

def add_audit_slide(prs):
    """Slide 26: Security Audit"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Security Audit Results"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = DEEP_BLUE
    
    content = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5))
    tf = content.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "8 CVE-Level Vulnerabilities Patched"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RED
    
    cves = [
        "🔴 HIGH: JWT Secret Exposure → Encrypted storage",
        "🔴 HIGH: SSRF via DNS Rebinding → Private IP blocking",
        "🔴 HIGH: Open Redirect → Strict URL validation",
        "🟡 MED: Path Traversal → Directory traversal detection",
        "🟡 MED: Information Disclosure → Sanitized errors",
        "🟡 MED: Weak Random Numbers → Secrets module",
        "🟢 LOW: Insecure Deserialization → Pydantic validation",
        "🟢 LOW: Missing Rate Limiting → 100 req/min limit",
    ]
    
    for cve in cves:
        p = tf.add_paragraph()
        p.text = cve
        p.font.size = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "\nBefore: ❌ Plaintext passwords, No authentication, No validation"
    p.font.size = Pt(14)
    p.font.color.rgb = RED
    
    p = tf.add_paragraph()
    p.text = "After: ✅ Fernet encrypted, JWT required, Enterprise validation"
    p.font.size = Pt(14)
    p.font.color.rgb = GREEN

def add_coverage_slide(prs):
    """Slide 27: Test Coverage"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Test Coverage & Quality"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = DEEP_BLUE
    
    content = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5))
    tf = content.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Coverage Metrics:"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = DEEP_BLUE
    
    covs = [
        "• Security-Critical Code: 100%",
        "• Feature Extraction: 85%",
        "• API Endpoints: 90%",
        "• Overall: 87%",
    ]
    
    for cov in covs:
        p = tf.add_paragraph()
        p.text = cov
        p.font.size = Pt(14)
    
    p = tf.add_paragraph()
    p.text = "\n14 Comprehensive Test Suites:"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = DEEP_BLUE
    
    p = tf.add_paragraph()
    p.text = "435 total test cases | 100% passing rate"
    p.font.size = Pt(16)
    p.font.color.rgb = GREEN

def add_testing_slide(prs):
    """Slide 28: Real-World Testing"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Real-World Testing Results"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = DEEP_BLUE
    
    content = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5))
    tf = content.text_frame
    tf.word_wrap = True
    
    tests = [
        ("Test 1: IDN Detection", "раураl.com → ✅ PHISHING (96.3%)", GREEN),
        ("Test 2: AI-Generated", "ChatGPT email → ✅ AI_GENERATED", GREEN),
        ("Test 3: Phishing Kit", "Gophish page → ✅ PHISHING_KIT", GREEN),
        ("Test 4: Browser Extension", "50 links → ✅ 47 safe, 3 phishing", GREEN),
    ]
    
    y_offset = 0
    for test_name, test_result, color in tests:
        p = tf.paragraphs[0] if y_offset == 0 else tf.add_paragraph()
        p.text = f"\n{test_name}"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = DEEP_BLUE
        
        p = tf.add_paragraph()
        p.text = test_result
        p.font.size = Pt(14)
        p.font.color.rgb = color
        
        y_offset += 1
    
    p = tf.add_paragraph()
    p.text = "\nField Test: 1,247 URLs | 98.4% True Positives | 0.6% False Positives"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = GREEN

def add_matrix_slide(prs):
    """Slide 29: Comparison Matrix"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Detailed Competitive Comparison"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = DEEP_BLUE
    
    content = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5))
    tf = content.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "21 Categories Compared"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = DEEP_BLUE
    
    p = tf.add_paragraph()
    p.text = "\nPhishing Guard leads in 17/21 categories!"
    p.font.size = Pt(16)
    p.font.color.rgb = GREEN
    
    highlights = [
        "\nKey Differentiators:",
        "• Only solution with IDN detection",
        "• Only solution with AI-generated classification",
        "• Only solution with 4-category output",
        "• Only open-source with enterprise security",
    ]
    
    for hl in highlights:
        p = tf.add_paragraph()
        p.text = hl
        p.font.size = Pt(14)

def add_uvp_slide(prs):
    """Slide 30: Unique Value Proposition"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Unique Value Proposition"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = DEEP_BLUE
    
    content = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5))
    tf = content.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = '"The only open-source phishing detection system that combines IDN protection, AI-generated detection, and enterprise security in a production-ready platform"'
    p.font.size = Pt(16)
    p.font.italic = True
    p.font.color.rgb = DEEP_BLUE
    
    p = tf.add_paragraph()
    p.text = "\n\n3 Pillars of Value:"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = DEEP_BLUE
    
    pillars = [
        "🥇 IDN FIRST - Only academic system with full Unicode",
        "🤖 AI DETECTION - Addresses ChatGPT threats",
        "🔒 ENTERPRISE SECURITY - 8 CVEs patched",
    ]
    
    for pillar in pillars:
        p = tf.add_paragraph()
        p.text = pillar
        p.font.size = Pt(14)
    
    p = tf.add_paragraph()
    p.text = "\n\nImpact: 99.8% F1 | 93 Features | 8 CVEs | Production-Ready"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = GOLD

def add_achievements_slide(prs):
    """Slide 31: Key Achievements"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Project Achievements"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = DEEP_BLUE
    
    achievements = [
        ("🏆", "Novel Research", "First IDN detection, AI classification, 4-tier architecture"),
        ("📊", "Exceptional Performance", "99.8% F1, 365% feature improvement, <150ms latency"),
        ("🔒", "Enterprise Security", "8 CVEs patched, production-ready hardening"),
        ("🌐", "Comprehensive Solution", "5 interfaces, Docker ready, MLflow integrated"),
        ("📖", "Open Source Impact", "10,775 LOC, 800+ docs, free for community"),
    ]
    
    y_pos = 1.2
    for emoji, ach_title, ach_desc in achievements:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y_pos), Inches(12), Inches(1))
        box.fill.solid()
        box.fill.fore_color.rgb = GOLD
        box.line.fill.background()
        
        text = slide.shapes.add_textbox(Inches(0.7), Inches(y_pos + 0.1), Inches(11.6), Inches(0.8))
        tf = text.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = f"{emoji} {ach_title}"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = DEEP_BLUE
        
        p = tf.add_paragraph()
        p.text = ach_desc
        p.font.size = Pt(12)
        p.font.color.rgb = BLACK
        
        y_pos += 1.15

def add_impact_slide(prs):
    """Slide 32: Impact"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Impact & Significance"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = DEEP_BLUE
    
    content = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5))
    tf = content.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Research Impact:"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = BLUE
    
    research = [
        "• Novel IDN detection methodology",
        "• First 4-category classification",
        "• Addresses AI-powered threats",
        "• Reproducible with MLflow",
    ]
    
    for r in research:
        p = tf.add_paragraph()
        p.text = r
        p.font.size = Pt(14)
    
    p = tf.add_paragraph()
    p.text = "\nPractical Impact:"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = GREEN
    
    practical = [
        "• Protects against 3.4B daily phishing emails",
        "• Open source = free for everyone",
        "• Enterprise-ready for deployment",
        "• Educational tool for courses",
    ]
    
    for p_text in practical:
        p = tf.add_paragraph()
        p.text = p_text
        p.font.size = Pt(14)

def add_limitations_slide(prs):
    """Slide 33: Limitations"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Limitations & Constraints"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = DEEP_BLUE
    
    limitations = [
        ("⚠️ MLLM Resources", "~2GB RAM required", "Optional tier, RF works standalone"),
        ("⚠️ Internet Dependency", "Tier 3-4 require connection", "Graceful degradation (99.8%→97.2%)"),
        ("⚠️ Training Data Bias", "English-centric", "IDN helps; multilingual planned"),
        ("⚠️ Browser Coverage", "Chrome/Brave/Edge only", "Firefox in roadmap"),
        ("⚠️ Email Scanner", "IMAP only", "Exchange API planned"),
    ]
    
    y_pos = 1.2
    for lim_title, lim_constraint, lim_mitigation in limitations:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y_pos), Inches(12), Inches(1))
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_RED
        box.line.color.rgb = RED
        
        text = slide.shapes.add_textbox(Inches(0.7), Inches(y_pos + 0.1), Inches(11.6), Inches(0.8))
        tf = text.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = lim_title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = RED
        
        p = tf.add_paragraph()
        p.text = f"Constraint: {lim_constraint} | Mitigation: {lim_mitigation}"
        p.font.size = Pt(12)
        p.font.color.rgb = BLACK
        
        y_pos += 1.15

def add_future_slide(prs):
    """Slide 34: Future Work"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Future Work & Roadmap"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = DEEP_BLUE
    
    content = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5))
    tf = content.text_frame
    tf.word_wrap = True
    
    phases = [
        ("Phase 3 (80%)", "Browser extension, Firefox support, Mobile app"),
        ("Phase 4 (25%)", "Kubernetes, monitoring, auto-scaling"),
        ("Phase 5 (0%)", "Fine-tuned LLM, federated learning, SIEM integration"),
    ]
    
    for phase, desc in phases:
        p = tf.paragraphs[0] if phase == "Phase 3 (80%)" else tf.add_paragraph()
        p.text = f"\n{phase}:"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = BLUE
        
        p = tf.add_paragraph()
        p.text = desc
        p.font.size = Pt(14)
    
    p = tf.add_paragraph()
    p.text = "\nTimeline: Q1 2025 → Q2 2025 → Q3-Q4 2025"
    p.font.size = Pt(14)
    p.font.color.rgb = GRAY

def add_thankyou_slide(prs):
    """Slide 35: Thank You"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DEEP_BLUE
    bg.line.fill.background()
    
    # Thank you
    thanks = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(1))
    tf = thanks.text_frame
    p = tf.paragraphs[0]
    p.text = "Thank You!"
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Project name
    proj = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(12), Inches(0.8))
    tf = proj.text_frame
    p = tf.paragraphs[0]
    p.text = "🛡️ Phishing Guard v2.0 🛡️"
    p.font.size = Pt(32)
    p.font.color.rgb = GOLD
    p.alignment = PP_ALIGN.CENTER
    
    # Tagline
    tag = slide.shapes.add_textbox(Inches(0.5), Inches(3.7), Inches(12), Inches(0.6))
    tf = tag.text_frame
    p = tf.paragraphs[0]
    p.text = "Protecting Users Against Next-Generation Phishing Attacks"
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(200, 200, 200)
    p.alignment = PP_ALIGN.CENTER
    
    # Resources
    res = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(12), Inches(1.5))
    tf = res.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Project Resources:"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = GOLD
    p.alignment = PP_ALIGN.CENTER
    
    p = tf.add_paragraph()
    p.text = "github.com/BandiAkarsh/phishing_detection_project"
    p.font.size = Pt(14)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    p = tf.add_paragraph()
    p.text = "\nQuestions & Discussion"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = GREEN
    p.alignment = PP_ALIGN.CENTER

if __name__ == "__main__":
    create_presentation()
